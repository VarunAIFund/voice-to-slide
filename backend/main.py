from fastapi import FastAPI, File, Form, UploadFile, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from starlette.requests import Request
import os
import uuid
import shutil
from pathlib import Path
import json
from typing import Dict, Optional
import base64
import io
import urllib.request
import logging
import time
import moviepy.editor as mp
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from dotenv import load_dotenv

# Load environment variables from .env file
load_dotenv()

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s | %(levelname)s | %(name)s | %(message)s",
)
logger = logging.getLogger("voice_to_slide")

app = FastAPI(title="Voice-to-Slide Generator", version="1.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:3000",
        "http://localhost:3001",
        "http://127.0.0.1:3000",
        "http://127.0.0.1:3001",
    ],
    allow_origin_regex=r"https?://(localhost|127\.0\.0\.1):\d+",
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# In-memory job tracking
jobs: Dict[str, dict] = {}

# Initialize OpenAI client
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

def ensure_openai_key():
    if not os.getenv("OPENAI_API_KEY"):
        raise HTTPException(
            status_code=500,
            detail="OPENAI_API_KEY is missing. Add it to backend/.env before generating slides.",
        )

def log_job(job_id: str, message: str, **fields):
    details = " ".join(f"{k}={v}" for k, v in fields.items())
    suffix = f" {details}" if details else ""
    logger.info(f"[job:{job_id}] {message}{suffix}")

# Create temp directories relative to this file so the server can be launched
# from any working directory (e.g. `uvicorn backend.main:app` from the project root).
_BASE_DIR = Path(__file__).parent
UPLOAD_DIR = _BASE_DIR / "temp" / "uploads"
OUTPUT_DIR = _BASE_DIR / "temp" / "outputs"
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
AUDIO_EXTENSIONS = {".mp3", ".wav", ".m4a", ".aac", ".ogg", ".flac", ".mp4"}

@app.middleware("http")
async def request_logger(request: Request, call_next):
    start = time.perf_counter()
    try:
        response = await call_next(request)
        elapsed_ms = int((time.perf_counter() - start) * 1000)
        logger.info(
            "request method=%s path=%s status=%s duration_ms=%s",
            request.method,
            request.url.path,
            response.status_code,
            elapsed_ms,
        )
        return response
    except Exception:
        elapsed_ms = int((time.perf_counter() - start) * 1000)
        logger.exception(
            "request_failed method=%s path=%s duration_ms=%s",
            request.method,
            request.url.path,
            elapsed_ms,
        )
        raise

# Define color themes for professional presentations
COLOR_THEMES = {
    "corporate_blue": {
        "name": "Corporate Blue",
        "background": RGBColor(240, 248, 255),  # Light blue background
        "primary": RGBColor(25, 118, 210),      # Professional blue
        "secondary": RGBColor(69, 90, 100),     # Dark gray
        "accent": RGBColor(0, 150, 136),        # Teal accent
        "text": RGBColor(33, 33, 33),           # Dark gray text
        "light_text": RGBColor(117, 117, 117)  # Light gray text
    },
    "modern_green": {
        "name": "Modern Green",
        "background": RGBColor(248, 255, 248),  # Light green background
        "primary": RGBColor(76, 175, 80),       # Modern green
        "secondary": RGBColor(55, 71, 79),      # Dark blue-gray
        "accent": RGBColor(255, 152, 0),        # Orange accent
        "text": RGBColor(33, 33, 33),           # Dark gray text
        "light_text": RGBColor(117, 117, 117)  # Light gray text
    },
    "elegant_purple": {
        "name": "Elegant Purple",
        "background": RGBColor(250, 245, 255),  # Light purple background
        "primary": RGBColor(156, 39, 176),      # Elegant purple
        "secondary": RGBColor(69, 39, 160),     # Deep purple
        "accent": RGBColor(255, 193, 7),        # Gold accent
        "text": RGBColor(33, 33, 33),           # Dark gray text
        "light_text": RGBColor(117, 117, 117)  # Light gray text
    },
    "professional_gray": {
        "name": "Professional Gray",
        "background": RGBColor(250, 250, 250),  # Light gray background
        "primary": RGBColor(96, 125, 139),      # Blue-gray
        "secondary": RGBColor(55, 71, 79),      # Dark gray
        "accent": RGBColor(255, 87, 34),        # Orange accent
        "text": RGBColor(33, 33, 33),           # Dark gray text
        "light_text": RGBColor(117, 117, 117)  # Light gray text
    }
}

def rgb_to_hex(rgb: RGBColor) -> str:
    # python-pptx RGBColor string form is a 6-char hex value like "1976D2"
    return f"#{str(rgb)}"

def extract_json_object(raw_text: str) -> dict:
    """Extract JSON object from model output, including fenced blocks."""
    if not raw_text:
        raise ValueError("Model returned empty slide payload")
    text = raw_text.strip()
    if text.startswith("```"):
        lines = [line for line in text.splitlines() if not line.strip().startswith("```")]
        text = "\n".join(lines).strip()
    start = text.find("{")
    end = text.rfind("}")
    if start == -1 or end == -1 or end <= start:
        raise ValueError("Slide payload did not contain a JSON object")
    return json.loads(text[start : end + 1])

def generate_slide_image(prompt: str, visual_style: str) -> Optional[bytes]:
    """Generate an image for a slide using OpenAI Images API."""
    enhanced_prompt = (
        f"Create a clean, presentation-ready hero image in {visual_style} style. "
        f"Use this concept: {prompt}. "
        "No text, no logos, no watermarks, high contrast, modern composition."
    )
    image = client.images.generate(
        model="gpt-image-1",
        prompt=enhanced_prompt,
        size="1536x1024",
    )
    if not image.data:
        return None
    first = image.data[0]
    if getattr(first, "b64_json", None):
        return base64.b64decode(first.b64_json)
    if getattr(first, "url", None):
        with urllib.request.urlopen(first.url, timeout=30) as response:
            return response.read()
    return None

def transcribe_audio_file(audio: UploadFile, job_id: str) -> str:
    audio_suffix = Path(audio.filename or "audio.wav").suffix.lower()
    if audio_suffix not in AUDIO_EXTENSIONS:
        raise HTTPException(status_code=400, detail="Unsupported audio format")

    audio_path = UPLOAD_DIR / f"{job_id}{audio_suffix or '.wav'}"
    with open(audio_path, "wb") as buffer:
        shutil.copyfileobj(audio.file, buffer)

    with open(audio_path, "rb") as audio_stream:
        result = client.audio.transcriptions.create(
            model="whisper-1",
            file=audio_stream,
        )
    return (result.text or "").strip()

def apply_slide_theme(slide, theme_colors, is_title_slide=False):
    """Apply base theme colors and typography to a slide."""
    # Set slide background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = theme_colors["background"]
    
    # Style title
    if slide.shapes.title:
        title_frame = slide.shapes.title.text_frame
        title_frame.clear()
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER if is_title_slide else PP_ALIGN.LEFT
        
        # Set title font properties
        title_font = title_para.font
        title_font.name = "Calibri"
        title_font.size = Pt(36) if is_title_slide else Pt(28)
        title_font.bold = True
        title_font.color.rgb = theme_colors["primary"]
    
    return slide

def add_slide_footer(slide, slide_number, total_slides, theme_colors, presentation_title):
    """Add footer with slide number and presentation title"""
    
    # Add slide number in bottom right
    slide_num_box = slide.shapes.add_textbox(
        left=Inches(8.5), 
        top=Inches(7), 
        width=Inches(1), 
        height=Inches(0.5)
    )
    slide_num_frame = slide_num_box.text_frame
    slide_num_frame.clear()
    slide_num_para = slide_num_frame.paragraphs[0]
    slide_num_para.text = f"{slide_number}/{total_slides}"
    slide_num_para.alignment = PP_ALIGN.RIGHT
    
    # Style slide number
    slide_num_font = slide_num_para.font
    slide_num_font.name = "Calibri"
    slide_num_font.size = Pt(12)
    slide_num_font.color.rgb = theme_colors["light_text"]
    
    # Add presentation title in bottom left (except for title slide)
    if slide_number > 1:
        title_box = slide.shapes.add_textbox(
            left=Inches(0.5), 
            top=Inches(7), 
            width=Inches(6), 
            height=Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_frame.clear()
        title_para = title_frame.paragraphs[0]
        title_para.text = presentation_title[:50] + "..." if len(presentation_title) > 50 else presentation_title
        title_para.alignment = PP_ALIGN.LEFT
        
        # Style footer title
        title_font = title_para.font
        title_font.name = "Calibri"
        title_font.size = Pt(10)
        title_font.color.rgb = theme_colors["light_text"]
    
    return slide

def add_decorative_elements(slide, theme_colors, is_title_slide=False):
    """Add decorative elements like lines and shapes"""
    
    if not is_title_slide:
        # Add a subtle line under the title
        line_shape = slide.shapes.add_connector(
            connector_type=1,  # Straight line
            begin_x=Inches(0.5), 
            begin_y=Inches(1.8), 
            end_x=Inches(9.5), 
            end_y=Inches(1.8)
        )
        
        # Style the line
        line = line_shape.line
        line.color.rgb = theme_colors["accent"]
        line.width = Pt(2)
    
    else:
        # Add decorative accent shape for title slide
        accent_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left=Inches(0.5),
            top=Inches(4.5),
            width=Inches(9),
            height=Inches(0.1)
        )
        
        # Style accent shape
        accent_fill = accent_shape.fill
        accent_fill.solid()
        accent_fill.fore_color.rgb = theme_colors["accent"]
        
        # Remove outline
        accent_shape.line.fill.background()
    
    return slide

def add_hero_layout(slide, slide_data, theme_colors, include_images, visual_style):
    """Create a split layout with bullets and optional generated image."""
    # Slide title
    title_box = slide.shapes.add_textbox(
        left=Inches(0.6),
        top=Inches(0.45),
        width=Inches(12.0),
        height=Inches(0.9),
    )
    title_frame = title_box.text_frame
    title_frame.clear()
    title_para = title_frame.paragraphs[0]
    title_para.text = slide_data["title"]
    title_para.alignment = PP_ALIGN.LEFT
    title_font = title_para.font
    title_font.name = "Calibri"
    title_font.size = Pt(34)
    title_font.bold = True
    title_font.color.rgb = theme_colors["primary"]

    # Left text panel
    text_panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left=Inches(0.55),
        top=Inches(1.45),
        width=Inches(6.35),
        height=Inches(5.25),
    )
    text_panel.fill.solid()
    text_panel.fill.fore_color.rgb = RGBColor(255, 255, 255)
    text_panel.fill.transparency = 6
    text_panel.line.fill.background()

    body_box = slide.shapes.add_textbox(
        left=Inches(0.95),
        top=Inches(1.85),
        width=Inches(5.6),
        height=Inches(4.6),
    )
    body_frame = body_box.text_frame
    body_frame.clear()

    for i, point in enumerate(slide_data.get("content", [])):
        para = body_frame.paragraphs[0] if i == 0 else body_frame.add_paragraph()
        para.text = point
        para.level = 0
        para.alignment = PP_ALIGN.LEFT
        para.space_after = Pt(8)
        font = para.font
        font.name = "Calibri"
        font.size = Pt(20)
        font.color.rgb = theme_colors["text"]

    # Right visual panel
    visual_panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left=Inches(7.1),
        top=Inches(1.45),
        width=Inches(5.65),
        height=Inches(5.25),
    )
    visual_panel.fill.solid()
    visual_panel.fill.fore_color.rgb = theme_colors["primary"]
    visual_panel.fill.transparency = 86
    visual_panel.line.fill.background()

    if include_images:
        prompt = slide_data.get("visual_prompt") or f"{slide_data['title']} cinematic concept art"
        image_bytes = generate_slide_image(prompt, visual_style)
        if image_bytes:
            slide.shapes.add_picture(
                io.BytesIO(image_bytes),
                left=Inches(7.25),
                top=Inches(1.6),
                width=Inches(5.35),
                height=Inches(4.95),
            )

@app.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    if not file.filename.endswith('.mp4'):
        raise HTTPException(status_code=400, detail="Only MP4 files are supported")
    
    job_id = str(uuid.uuid4())
    file_path = UPLOAD_DIR / f"{job_id}.mp4"
    
    with open(file_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    
    jobs[job_id] = {
        "status": "uploaded",
        "filename": file.filename,
        "file_path": str(file_path),
        "progress": 10
    }
    
    log_job(job_id, "file_uploaded", filename=file.filename)
    return {"job_id": job_id, "message": "File uploaded successfully"}

@app.get("/status/{job_id}")
async def get_status(job_id: str):
    if job_id not in jobs:
        raise HTTPException(status_code=404, detail="Job not found")
    
    return jobs[job_id]

@app.get("/transcript/{job_id}")
async def get_transcript(job_id: str):
    if job_id not in jobs:
        raise HTTPException(status_code=404, detail="Job not found")
    
    job = jobs[job_id]
    
    if job["status"] == "uploaded":
        # Extract audio
        jobs[job_id]["status"] = "extracting_audio"
        jobs[job_id]["progress"] = 20
        log_job(job_id, "extracting_audio_started")
        
        try:
            video_path = job["file_path"]
            audio_path = UPLOAD_DIR / f"{job_id}.wav"
            
            video = mp.VideoFileClip(video_path)
            video.audio.write_audiofile(str(audio_path), verbose=False, logger=None)
            video.close()
            
            jobs[job_id]["status"] = "transcribing"
            jobs[job_id]["progress"] = 40
            jobs[job_id]["audio_path"] = str(audio_path)
            log_job(job_id, "transcription_started", audio_path=audio_path.name)
            
            # Transcribe audio via OpenAI API for broader Python compatibility
            with open(audio_path, "rb") as audio_stream:
                result = client.audio.transcriptions.create(
                    model="whisper-1",
                    file=audio_stream,
                )
            transcript = (result.text or "").strip()
            
            jobs[job_id]["status"] = "transcript_ready"
            jobs[job_id]["progress"] = 60
            jobs[job_id]["transcript"] = transcript
            
            log_job(job_id, "transcription_completed", transcript_chars=len(transcript))
            
        except Exception as e:
            jobs[job_id]["status"] = "error"
            jobs[job_id]["error"] = str(e)
            logger.exception("[job:%s] transcript_failed error=%s", job_id, str(e))
            raise HTTPException(status_code=500, detail=str(e))
    
    if "transcript" in jobs[job_id]:
        return {"transcript": jobs[job_id]["transcript"]}
    else:
        return {"message": "Transcript not ready yet"}

@app.post("/generate-slides/{job_id}")
async def generate_slides(
    job_id: str,
    theme: str = "corporate_blue",
    include_images: bool = True,
    visual_style: str = "cinematic",
):
    ensure_openai_key()
    if job_id not in jobs:
        raise HTTPException(status_code=404, detail="Job not found")
    
    job = jobs[job_id]
    
    if "transcript" not in job:
        raise HTTPException(status_code=400, detail="Transcript not available")
    
    try:
        jobs[job_id]["status"] = "generating_slides"
        jobs[job_id]["progress"] = 70
        log_job(job_id, "slide_generation_started", theme=theme, include_images=include_images, visual_style=visual_style)
        
        # Generate slide content using OpenAI
        transcript = job["transcript"]
        
        prompt = f"""
        Convert this transcript into a polished presentation outline.
        Requirements:
        - Create 5-8 slides
        - Keep each slide concise: 3-5 bullets, each 6-16 words
        - Bullets should be actionable and not repetitive
        - Provide a visual concept prompt for each slide
        - Return JSON only with no markdown

        JSON schema:
        {{
          "title": "Presentation Title",
          "subtitle": "Optional subtitle",
          "slides": [
            {{
              "title": "Slide title",
              "content": ["Bullet 1", "Bullet 2", "Bullet 3"],
              "visual_prompt": "A concise visual concept for AI image generation"
            }}
          ]
        }}

        Transcript:
        {transcript}
        """
        
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}],
            max_tokens=1800,
            response_format={"type": "json_object"},
        )
        slide_content = extract_json_object(response.choices[0].message.content)
        
        jobs[job_id]["status"] = "creating_powerpoint"
        jobs[job_id]["progress"] = 85
        jobs[job_id]["slide_content"] = slide_content
        log_job(job_id, "slide_outline_ready", slides=len(slide_content.get("slides", [])))
        
        # Get selected theme colors
        theme_colors = COLOR_THEMES.get(theme, COLOR_THEMES["corporate_blue"])
        
        # Create a widescreen PowerPoint presentation
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Calculate total slides
        total_slides = len(slide_content["slides"]) + 1
        presentation_title = slide_content["title"]
        
        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_slide_theme(title_slide, theme_colors, is_title_slide=True)
        title_band = title_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.15)
        )
        title_band.fill.solid()
        title_band.fill.fore_color.rgb = theme_colors["primary"]
        title_band.line.fill.background()

        title_box = title_slide.shapes.add_textbox(
            left=Inches(0.75),
            top=Inches(1.65),
            width=Inches(11.8),
            height=Inches(2.25),
        )
        title_frame = title_box.text_frame
        title_frame.clear()
        title_para = title_frame.paragraphs[0]
        title_para.text = presentation_title
        title_para.alignment = PP_ALIGN.LEFT
        title_font = title_para.font
        title_font.name = "Calibri"
        title_font.size = Pt(50)
        title_font.bold = True
        title_font.color.rgb = theme_colors["secondary"]

        subtitle_text = slide_content.get("subtitle") or "Generated from audio transcript"
        subtitle_box = title_slide.shapes.add_textbox(
            left=Inches(0.8),
            top=Inches(4.45),
            width=Inches(8.5),
            height=Inches(0.8),
        )
        subtitle_para = subtitle_box.text_frame.paragraphs[0]
        subtitle_para.text = subtitle_text
        subtitle_para.font.name = "Calibri"
        subtitle_para.font.size = Pt(20)
        subtitle_para.font.color.rgb = theme_colors["light_text"]
        add_slide_footer(title_slide, 1, total_slides, theme_colors, presentation_title)

        # Content slides
        for slide_index, slide_data in enumerate(slide_content["slides"]):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            apply_slide_theme(slide, theme_colors, is_title_slide=False)
            add_hero_layout(
                slide=slide,
                slide_data=slide_data,
                theme_colors=theme_colors,
                include_images=include_images,
                visual_style=visual_style,
            )
            add_slide_footer(slide, slide_index + 2, total_slides, theme_colors, presentation_title)
        
        # Save PowerPoint file
        ppt_path = OUTPUT_DIR / f"{job_id}.pptx"
        prs.save(str(ppt_path))
        
        jobs[job_id]["status"] = "completed"
        jobs[job_id]["progress"] = 100
        jobs[job_id]["ppt_path"] = str(ppt_path)
        
        log_job(job_id, "powerpoint_generated", output=ppt_path.name)
        
        return {"message": "Slides generated successfully", "slide_content": slide_content}
        
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        logger.exception("[job:%s] slide_generation_failed error=%s", job_id, str(e))
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/generate-from-input")
async def generate_from_input(
    transcript: str = Form(default=""),
    theme: str = Form(default="corporate_blue"),
    include_images: bool = Form(default=True),
    visual_style: str = Form(default="cinematic"),
    audio: Optional[UploadFile] = File(default=None),
):
    text_input = (transcript or "").strip()
    audio_transcript = ""
    job_id = str(uuid.uuid4())

    if audio:
        logger.info("generate_from_input audio_received filename=%s", audio.filename)
        audio_transcript = transcribe_audio_file(audio, job_id)

    final_transcript = "\n\n".join(part for part in [text_input, audio_transcript] if part).strip()
    if not final_transcript:
        raise HTTPException(status_code=400, detail="Provide text input or an audio file")

    jobs[job_id] = {
        "status": "transcript_ready",
        "filename": (audio.filename if audio and audio.filename else "text-input"),
        "progress": 60,
        "transcript": final_transcript,
    }

    result = await generate_slides(
        job_id=job_id,
        theme=theme,
        include_images=include_images,
        visual_style=visual_style,
    )
    result["job_id"] = job_id
    result["transcript"] = final_transcript
    log_job(job_id, "generate_from_input_completed")
    return result

@app.get("/download/{job_id}")
async def download_file(job_id: str):
    if job_id not in jobs:
        raise HTTPException(status_code=404, detail="Job not found")
    
    job = jobs[job_id]
    
    if job["status"] != "completed" or "ppt_path" not in job:
        raise HTTPException(status_code=400, detail="PowerPoint not ready")
    
    ppt_path = job["ppt_path"]
    
    if not os.path.exists(ppt_path):
        raise HTTPException(status_code=404, detail="File not found")
    
    return FileResponse(
        path=ppt_path,
        filename=f"{job['filename']}.pptx",
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

@app.get("/themes")
async def get_themes():
    """Get available presentation themes"""
    themes = {}
    for theme_id, theme_data in COLOR_THEMES.items():
        themes[theme_id] = {
            "name": theme_data["name"],
            "colors": {
                "primary": rgb_to_hex(theme_data["primary"]),
                "secondary": rgb_to_hex(theme_data["secondary"]),
                "accent": rgb_to_hex(theme_data["accent"]),
                "background": rgb_to_hex(theme_data["background"]),
            }
        }
    return {"themes": themes}

@app.get("/")
async def root():
    return {"message": "Voice-to-Slide Generator API"}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)